#!/usr/bin/env python3
"""Run an end-to-end smoke test for the professor skill workspace flow."""

from __future__ import annotations

import json
import shutil
import subprocess
import sys
from collections import Counter
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
TMP = ROOT / "temp_smoke"


def run(*args: str) -> str:
    result = subprocess.run(
        [sys.executable, *args],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=True,
    )
    return result.stdout.strip()


def write_file(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def make_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_heading("数据库系统作业", level=1)
    document.add_paragraph("定义和场景都要写，不要只写缩写。")
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def make_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "事务与并发控制"
    slide.placeholders[1].text = "ACID\n锁协议\n可串行化\n这一步每年都有人丢分"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def make_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    path.parent.mkdir(parents=True, exist_ok=True)
    pdf = canvas.Canvas(str(path), pagesize=A4)
    pdf.setFont("Helvetica", 12)
    lines = [
        "Final sample",
        "1. Explain what ACID means.",
        "2. Analyze why two phase locking gives serializability.",
        "Answers without reasoning will not receive full credit.",
    ]
    y = 780
    for line in lines:
        pdf.drawString(72, y, line)
        y -= 24
    pdf.save()


def main() -> int:
    if TMP.exists():
        shutil.rmtree(TMP)
    TMP.mkdir(parents=True, exist_ok=True)

    workspace = run(
        "tools/professor_writer.py",
        "--name",
        "测试老师",
        "--course",
        "数据库系统",
        "--school",
        "某大学",
        "--department",
        "计算机学院",
        "--base-dir",
        str(TMP),
    )
    professor_dir = Path(workspace)

    make_pdf(professor_dir / "materials" / "exams" / "sample-final.pdf")
    make_docx(professor_dir / "materials" / "assignments" / "hw1.docx")
    make_pptx(professor_dir / "materials" / "slides" / "week08.pptx")
    write_file(
        professor_dir / "materials" / "transcripts" / "lecture-01.txt",
        "这个定义不要背缩写，要知道它为什么这样设计。期末如果只写名词解释，分数不会高。\n",
    )
    write_file(
        professor_dir / "materials" / "chats" / "qa.txt",
        "学生：老师这个会考吗？\n老师：作业最后一题先看懂。\n",
    )

    run("tools/build_professor_outputs.py", str(professor_dir), "--strict")

    output = run("tools/validate_professor.py", str(professor_dir), "--strict")
    if "STATUS: valid" not in output:
        raise SystemExit(output)

    for filename in ("persona.md", "course.md", "review_guide.md"):
        content = (professor_dir / filename).read_text(encoding="utf-8")
        if "[fill me]" in content:
            raise SystemExit(f"{filename} still has placeholders")

    index = json.loads((professor_dir / "exports" / "extracted" / "index.json").read_text(encoding="utf-8"))
    statuses = Counter(record["extract_status"] for record in index)
    if statuses.get("ok", 0) < 5:
        raise SystemExit(f"expected parsed files, got statuses={dict(statuses)}")

    shutil.rmtree(TMP)
    print("smoke test passed")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

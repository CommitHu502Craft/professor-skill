#!/usr/bin/env python3
"""Extract text from professor source materials into a unified text layer."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from materials_lib import (
    EXTRACTABLE_SUFFIXES,
    MATERIAL_DIRS,
    extracted_relative_path,
    infer_category,
    normalize_text,
    read_text_best_effort,
    save_extract_index,
)


def extract_text_file(source: Path) -> tuple[str, str]:
    text = normalize_text(read_text_best_effort(source))
    return text, "text"


def extract_pdf_file(source: Path) -> tuple[str, str]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - dependency controlled by requirements
        raise RuntimeError("missing dependency: pypdf") from exc

    reader = PdfReader(str(source))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = normalize_text(page.extract_text() or "")
        if text:
            pages.append(f"[Page {index}]\n{text}")
    return "\n\n".join(pages).strip(), "pypdf"


def extract_docx_file(source: Path) -> tuple[str, str]:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("missing dependency: python-docx") from exc

    document = Document(str(source))
    paragraphs = [normalize_text(paragraph.text) for paragraph in document.paragraphs]
    text = "\n".join(paragraph for paragraph in paragraphs if paragraph)
    return text.strip(), "python-docx"


def extract_pptx_file(source: Path) -> tuple[str, str]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("missing dependency: python-pptx") from exc

    presentation = Presentation(str(source))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = normalize_text(shape.text)
                if text:
                    chunks.append(text)
        if chunks:
            slides.append(f"[Slide {index}]\n" + "\n".join(chunks))
    return "\n\n".join(slides).strip(), "python-pptx"


def extract_file(source: Path) -> tuple[str, str]:
    suffix = source.suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json", ".srt"}:
        return extract_text_file(source)
    if suffix == ".pdf":
        return extract_pdf_file(source)
    if suffix == ".docx":
        return extract_docx_file(source)
    if suffix == ".pptx":
        return extract_pptx_file(source)
    return "", "unsupported"


def update_meta(root: Path, records: list[dict[str, object]]) -> None:
    meta_path = root / "meta.json"
    if not meta_path.exists():
        return
    meta = json.loads(meta_path.read_text(encoding="utf-8"))
    meta["sources"] = records
    meta.setdefault("artifacts", {})
    meta["artifacts"]["extraction_index"] = "exports/extracted/index.json"
    meta_path.write_text(json.dumps(meta, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("professor_dir", help="Professor workspace directory")
    parser.add_argument("--update-meta", action="store_true", help="Write extraction records into meta.json")
    args = parser.parse_args(argv)

    root = Path(args.professor_dir).resolve()
    materials_root = root / "materials"
    if not root.exists():
        raise SystemExit(f"professor workspace not found: {root}")
    if not materials_root.exists():
        raise SystemExit(f"materials directory not found: {materials_root}")

    files = sorted(
        path for path in materials_root.rglob("*") if path.is_file() and not path.name.startswith(".")
    )
    records: list[dict[str, object]] = []

    for source_file in files:
        category = infer_category(source_file)
        relative = source_file.relative_to(root).as_posix()
        extract_record = {
            "path": relative,
            "category": category,
            "bytes": source_file.stat().st_size,
            "extract_status": "skipped",
            "extracted_path": None,
            "text_length": 0,
            "extractor": None,
        }

        if source_file.suffix.lower() not in EXTRACTABLE_SUFFIXES:
            extract_record["extract_status"] = "unsupported"
            records.append(extract_record)
            continue

        try:
            text, extractor = extract_file(source_file)
        except Exception as exc:  # pragma: no cover - tested via script behavior
            extract_record["extract_status"] = "error"
            extract_record["error"] = str(exc)
            records.append(extract_record)
            continue

        if not text:
            extract_record["extract_status"] = "empty"
            extract_record["extractor"] = extractor
            records.append(extract_record)
            continue

        target = root / extracted_relative_path(materials_root, source_file)
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(text + "\n", encoding="utf-8")

        extract_record["extract_status"] = "ok"
        extract_record["extracted_path"] = target.relative_to(root).as_posix()
        extract_record["text_length"] = len(text)
        extract_record["extractor"] = extractor
        records.append(extract_record)

    save_extract_index(root, records)
    if args.update_meta:
        update_meta(root, records)

    summary = {
        "total_files": len(records),
        "ok": sum(1 for record in records if record["extract_status"] == "ok"),
        "unsupported": sum(1 for record in records if record["extract_status"] == "unsupported"),
        "empty": sum(1 for record in records if record["extract_status"] == "empty"),
        "error": sum(1 for record in records if record["extract_status"] == "error"),
    }
    print(json.dumps(summary, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
